
from __future__ import annotations

import argparse
import os
import sys
import traceback
import json
from dataclasses import dataclass
from pathlib import Path
from datetime import datetime
from zipfile import ZipFile
import xml.etree.ElementTree as ET
from typing import Optional, Union

if sys.stdout and hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(errors="replace")
if sys.stderr and hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(errors="replace")

from .dataplotter import DataPlotter
from .logger import log, configure as configure_logging

DEFAULT_FIG_SIZE = {
    "waveform": (15.5, 6.4),
    "scatter": (10, 8),
    "psd": (10, 8),
    "histogram": (10, 8),
    "bar": (10, 6),
}
_ASPECT_RATIO_CACHE = {}

@dataclass
class PlotJobConfig:

    title: str
    root_folder: Path
    output_dir: Path
    runs: list
    plot_definitions: tuple
    channel_mappings: Optional[dict] = None
    channel_transforms: Optional[dict] = None
    calculated_channels: Optional[dict] = None
    filters: Optional[dict] = None
    resample_rate: Optional[float] = None
    units_map: Optional[dict] = None
    fig_size: Optional[Union[list, dict]] = None
    scatter_max_points: int = 45000
    bar_secondary_axis_ratio: float = 20.0
    box_plot_settings: Optional[dict] = None
    verbose: bool = False
    powerpoint_template: Optional[Path] = None
    powerpoint_output: Optional[Path] = None
    export_map: Optional[list] = None
    powerpoint_start_slide: int = 1
    powerpoint_exports: Optional[list] = None  # list of (template, output, export_map[, start_slide]) tuples
    open_output: bool = True
    output_dpi: int = 300
    vibrations_fit: Optional[dict] = None

def Slide(layout: str, *plot_refs: str) -> dict:
    images = [_plot_ref_to_filename(ref) for ref in plot_refs]
    return {"layout": layout, "images": images}

def _plot_ref_to_filename(ref: str) -> str:
    if "/" in ref:
        prefix, name = ref.split("/", 1)
    else:
        return ref if ref.endswith(".png") else f"{ref}.png"
    safe = (
        name.lower()
        .replace(" ", "_")
        .replace("(", "").replace(")", "")
        .replace("/", "_").replace("\\", "_")
    )
    pref = prefix.lower()
    return f"{pref}/{pref}_{safe}.png"

def _resolve_export_map(export_map, plot_definitions, start_slide=1):
    if export_map is None:
        return None
    offset = max(1, int(start_slide))
    return {i + offset: slide for i, slide in enumerate(export_map)}

_VALID_RUN_TYPES = {"OC", "CAR", "DLS", "DIL"}
_VALID_RUN_FILETYPES = {".csv", ".parquet", ".txt"}
_FOLDER_RUN_COLOR_PALETTE = (
    "#FF8000", "#2000BF", "#D70000", "#008CFF",
    "#00CC88", "#CC0066", "#FFD700", "#4C00BF",
)
_TYPE_COLORMAPS = {
    "CAR": "Oranges",
    "DLS": "Blues",
    "DIL": "Greens",
    "OC":  "Purples",
}

def _shades_from_cmap(
    cmap_name: str, n: int, offset: int = 0,
    low: float = 0.45, high: float = 0.95,
) -> list[str]:
    import matplotlib.cm as _cm
    from matplotlib.colors import to_hex
    cmap = _cm.get_cmap(cmap_name)
    total = max(n + offset, 2)
    span = high - low
    pts = [low + span * (i / max(total - 1, 1)) for i in range(offset, offset + n)]
    return [to_hex(cmap(p)) for p in pts]

def _interleave_indices(n: int) -> list[int]:
    """Return a permutation of range(n) that alternates between the
    two ends of the sequence (0, n-1, 1, n-2, 2, n-3, ...).

    Used to colour split-by children so neighbouring legend entries land
    at opposite ends of the colormap and contrast strongly, while the
    full set still spans the type's hue range.
    """
    if n <= 1:
        return list(range(n))
    lo, hi = 0, n - 1
    out: list[int] = []
    while lo <= hi:
        out.append(lo)
        lo += 1
        if lo <= hi:
            out.append(hi)
            hi -= 1
    return out

def _interpolate_two_colors(start: str, end: str, n: int) -> list[str]:
    """Return ``n`` colours interpolated through HSV space between two
    user-supplied endpoints (hex strings, named colours, or any
    matplotlib-recognised colour spec).

    Hue is interpolated linearly without shortest-arc wrap correction so
    the path traces the rainbow when the endpoints span the spectrum
    (e.g. red ``#FF0000`` to blue-violet ``#4800FF`` passes through
    orange, yellow, green, cyan, blue). Saturation and value are also
    interpolated linearly so the endpoints reproduce the user's colours
    exactly.

    For ``n == 1`` returns the midpoint; for ``n >= 2`` the first
    colour is ``start`` and the last is ``end``.
    """
    from matplotlib.colors import to_rgb, to_hex, rgb_to_hsv, hsv_to_rgb
    rgb1 = to_rgb(start)
    rgb2 = to_rgb(end)
    h1, s1, v1 = rgb_to_hsv(rgb1)
    h2, s2, v2 = rgb_to_hsv(rgb2)
    if n <= 0:
        return []
    if n == 1:
        mid_hsv = (
            ((h1 + h2) / 2) % 1.0,
            (s1 + s2) / 2,
            (v1 + v2) / 2,
        )
        return [to_hex(hsv_to_rgb(mid_hsv))]
    out: list[str] = []
    for i in range(n):
        t = i / (n - 1)
        h = (h1 + (h2 - h1) * t) % 1.0
        s = s1 + (s2 - s1) * t
        v = v1 + (v2 - v1) * t
        out.append(to_hex(hsv_to_rgb((h, s, v))))
    return out

_VALID_CONSOLIDATE_MODES = {True, "only"}

# Chronological session priority for F1 regular and sprint weekends. Lower is
# earlier. Unknown tokens sort after known ones (alphabetical tiebreak). Used
# by folder consolidation grouping and by the modal-evolution compare layout
# so RED vs BLUE line up in event order (not alphabetical).
SESSION_PRIORITY = {
    "P1": 0, "FP1": 0,
    "P2": 1, "FP2": 1,
    "P3": 2, "FP3": 2,
    "SQ": 3,
    "SR": 4,
    "Q":  5,
    "GP": 6,
}

def _session_sort_key(token: str) -> tuple:
    """Sort key for session tokens: known → (priority, token), unknown →
    (len(SESSION_PRIORITY), token). Preserves chronological order for known
    tokens and gives a stable alphabetical fallback for anything else.
    """
    tok = str(token).upper()
    return (SESSION_PRIORITY.get(tok, len(SESSION_PRIORITY)), tok)

# Preset extractors for `consolidate_by`. Each takes a Path and returns the
# group key string (or None to exclude the file from any consolidated group).
def _session_token_from_stem(path: Path) -> Optional[str]:
    """Standard race-data filename convention places the session token at
    position -2 (e.g. ``<event>_<date>_<car>_<driver>_<session>_<run>``).
    Returns e.g. ``"P1"`` / ``"SQ"`` / ``"SR"`` / ``"Q"`` / ``"GP"``.
    """
    parts = path.stem.split("_")
    return parts[-2] if len(parts) >= 2 else None

_CONSOLIDATE_BY_PRESETS = {
    "session": _session_token_from_stem,
}

def _resolve_consolidate_by(spec):
    """Return a callable ``(Path) -> Optional[str]`` for a ``consolidate_by`` spec.

    Accepted forms:
      * ``None``           — no partitioning (all sources → single group).
      * ``"session"`` (str)— preset extractor for race-data filenames.
      * regex ``str``      — pattern with one capture group, applied to the
                             file stem; the captured group becomes the key.
      * ``callable``       — receives the file ``Path`` and returns the key.
    """
    if spec is None:
        return None
    if callable(spec):
        return spec
    if isinstance(spec, str):
        if spec in _CONSOLIDATE_BY_PRESETS:
            return _CONSOLIDATE_BY_PRESETS[spec]
        import re
        try:
            pattern = re.compile(spec)
        except re.error as exc:
            raise ValueError(
                f"consolidate_by={spec!r} is not a valid regex: {exc}"
            ) from exc
        if pattern.groups < 1:
            raise ValueError(
                f"consolidate_by regex {spec!r} must contain at least one "
                "capture group identifying the partition key."
            )
        def _extract(p: Path, _pat=pattern) -> Optional[str]:
            m = _pat.search(p.stem)
            return m.group(1) if m else None
        return _extract
    raise TypeError(
        f"consolidate_by must be None, a preset name, a regex string, or a "
        f"callable; got {type(spec).__name__}."
    )

def _make_consolidated_entry(
    source_entries: list,
    run: dict,
    ext: str,
    *,
    group_key: Optional[str] = None,
) -> dict:
    """Build a synthetic 'consolidated' run dict from already-expanded sources.

    The returned dict has no 'file' key but carries `_consolidate_sources`
    (list of source run names, lowercased) for DataPlotter to merge after
    preprocessing.
    """
    name_tpl = run.get("consolidated_name")
    if name_tpl and "{group}" in name_tpl:
        name = name_tpl.format(group=group_key or "consolidated")
    elif name_tpl and group_key:
        name = f"{name_tpl}_{group_key}"
    elif name_tpl:
        name = name_tpl
    elif run.get("name"):
        name = run["name"]
        if group_key:
            name = f"{name}_{group_key}"
    else:
        prefix = run.get("name_prefix", "")
        folder_label = run.get("folder", "ALL")
        if folder_label == ".":
            folder_label = "ALL"
        suffix = group_key or "consolidated"
        name = f"{prefix}{folder_label}_{suffix}".strip("_")
    color = run.get("consolidated_color") or run.get("color")
    if not color:
        run_type = run.get("type")
        if run_type and run_type in _TYPE_COLORMAPS:
            color = _shades_from_cmap(_TYPE_COLORMAPS[run_type], 1, offset=0)[0]
    entry = {
        "name": name,
        "_consolidate_sources": [s["name"].lower() for s in source_entries],
        "_consolidate_ext": ext,
    }
    if group_key:
        entry["_consolidate_group"] = group_key
    if color:
        entry["color"] = color
    for k in ("type", "group", "use_python_engine"):
        if k in run:
            entry[k] = run[k]
    for extra in ("session_index", "session_label"):
        if extra in run:
            entry[extra] = run[extra]
    return entry

def _expand_folder_runs(runs: list, root_folder: Path) -> list:
    if not runs:
        return list(runs)
    root_path = Path(root_folder).resolve()
    expanded: list = []
    type_offsets: dict[str, int] = {}
    for i, run in enumerate(runs):
        if "folder" not in run:
            if run.get("split_by") and (
                run.get("consolidate") or run.get("consolidate_by")
                or "_consolidate_sources" in run
            ):
                raise ValueError(
                    f"Run[{i}] {run.get('name', '?')!r}: split_by cannot be "
                    f"combined with consolidate / consolidate_by on the same entry."
                )
            expanded.append(run)
            continue
        folder_rel = run["folder"]
        filetype = run.get("filetype")
        if not filetype:
            raise ValueError(
                f"Folder run[{i}] {folder_rel!r}: missing 'filetype' "
                f"(one of {sorted(_VALID_RUN_FILETYPES)})."
            )
        ext = filetype.lower()
        if not ext.startswith("."):
            ext = "." + ext
        if ext not in _VALID_RUN_FILETYPES:
            raise ValueError(
                f"Folder run[{i}] {folder_rel!r}: filetype {filetype!r} "
                f"must be one of {sorted(_VALID_RUN_FILETYPES)}."
            )
        run_type = run.get("type")
        if run_type and run_type not in _VALID_RUN_TYPES:
            raise ValueError(
                f"Folder run[{i}] {folder_rel!r}: unknown type {run_type!r}. "
                f"Expected one of: {', '.join(sorted(_VALID_RUN_TYPES))}"
            )
        folder_path = (root_path / folder_rel).resolve()
        if not folder_path.is_dir():
            raise FileNotFoundError(
                f"Folder run[{i}]: directory not found -> {folder_path}"
            )
        files = sorted(
            p for p in folder_path.iterdir()
            if p.is_file() and p.suffix.lower() == ext
        )
        if not files:
            raise FileNotFoundError(
                f"Folder run[{i}] {folder_rel!r}: no '{ext}' files in {folder_path}."
            )
        contains = run.get("contains")
        if contains:
            if not isinstance(contains, str):
                raise ValueError(
                    f"Folder run[{i}] {folder_rel!r}: 'contains' must be a string, "
                    f"got {type(contains).__name__}."
                )
            needle = contains.lower()
            files = [p for p in files if needle in p.name.lower()]
            if not files:
                raise FileNotFoundError(
                    f"Folder run[{i}] {folder_rel!r}: no '{ext}' files in "
                    f"{folder_path} matched contains={contains!r}."
                )
        colors = run.get("colors")
        single_color = run.get("color")
        name_prefix = run.get("name_prefix", "")
        reserved = {"folder", "filetype", "colors", "name_prefix",
                    "name", "file", "color", "contains",
                    "consolidate", "consolidate_by",
                    "consolidated_name", "consolidated_color"}
        common = {k: v for k, v in run.items() if k not in reserved}
        n_files = len(files)
        if not colors and not single_color and run_type in _TYPE_COLORMAPS:
            offset = type_offsets.get(run_type, 0)
            auto_colors = _shades_from_cmap(_TYPE_COLORMAPS[run_type], n_files, offset)
            type_offsets[run_type] = offset + n_files
        else:
            auto_colors = None
        consolidate = run.get("consolidate")
        if consolidate not in (None, *_VALID_CONSOLIDATE_MODES):
            raise ValueError(
                f"Folder run[{i}] {folder_rel!r}: consolidate must be True or 'only'; "
                f"got {consolidate!r}."
            )
        if run.get("split_by") and consolidate is not None:
            raise ValueError(
                f"Folder run[{i}] {folder_rel!r}: split_by cannot be combined "
                f"with consolidate / consolidate_by on the same entry."
            )
        try:
            partition_fn = _resolve_consolidate_by(run.get("consolidate_by"))
        except (ValueError, TypeError) as exc:
            raise ValueError(
                f"Folder run[{i}] {folder_rel!r}: {exc}"
            ) from exc
        produced: list = []
        produced_paths: list[Path] = []
        for j, f in enumerate(files):
            entry = dict(common)
            entry["name"] = f"{name_prefix}{f.stem}"
            try:
                entry["file"] = str(f.relative_to(root_path))
            except ValueError:
                entry["file"] = str(f)
            if colors:
                entry["color"] = colors[j % len(colors)]
            elif single_color:
                entry["color"] = single_color
            elif auto_colors is not None:
                entry["color"] = auto_colors[j]
            else:
                entry["color"] = _FOLDER_RUN_COLOR_PALETTE[
                    j % len(_FOLDER_RUN_COLOR_PALETTE)
                ]
            produced.append(entry)
            produced_paths.append(f)
        if consolidate is None or consolidate is True:
            expanded.extend(produced)
        if consolidate in (True, "only"):
            if partition_fn is None:
                cons = _make_consolidated_entry(produced, run, ext)
                if consolidate == "only":
                    cons["_consolidate_drop_sources"] = True
                    expanded.extend(produced)
                expanded.append(cons)
            else:
                groups: dict[str, list] = {}
                for entry, path in zip(produced, produced_paths):
                    key = partition_fn(path)
                    if key is None:
                        continue
                    groups.setdefault(str(key), []).append(entry)
                if not groups:
                    raise ValueError(
                        f"Folder run[{i}] {folder_rel!r}: consolidate_by produced "
                        f"no groups (all files returned None)."
                    )
                if consolidate == "only":
                    expanded.extend(produced)
                for gkey in sorted(groups.keys(), key=_session_sort_key):
                    gentries = groups[gkey]
                    cons = _make_consolidated_entry(
                        gentries, run, ext, group_key=gkey,
                    )
                    if consolidate == "only":
                        cons["_consolidate_drop_sources"] = True
                    expanded.append(cons)
                log.info(
                    "consolidate_by partitioned %d source file(s) into %d group(s): %s",
                    len(produced), len(groups),
                    ", ".join(sorted(groups, key=_session_sort_key)),
                )
    return expanded

def validate_config(config: PlotJobConfig) -> list[str]:
    issues: list[str] = []
    if not config.runs:
        issues.append("RUNS list is empty — nothing to plot.")
    # Run names are the primary key into run_data / run_sample_rates / modal_results.
    # Duplicates silently overwrite earlier loads and cause the preprocessing
    # pipeline (mappings, transforms, calculated channels, filters) to iterate
    # the same DataFrame twice, producing duplicate columns and a TypeError in
    # apply_filters. Catch it here so the message is actionable.
    name_to_indices: dict[str, list[int]] = {}
    for i, run in enumerate(config.runs):
        key = str(run.get("name", "")).strip().lower()
        if key:
            name_to_indices.setdefault(key, []).append(i)
    for key, idxs in name_to_indices.items():
        if len(idxs) > 1:
            issues.append(
                f"Duplicate run name '{config.runs[idxs[0]]['name']}' "
                f"at RUNS positions {idxs}. Run names must be unique — "
                f"distinguish them (e.g. 'DLS RED R' vs 'DLS BLUE R')."
            )
    for i, run in enumerate(config.runs):
        label = run.get("name", f"<unnamed run[{i}]>")
        if not run.get("name"):
            issues.append(f"Run[{i}]: missing 'name' key.")
        is_consolidated = "_consolidate_sources" in run
        if is_consolidated:
            sources = run.get("_consolidate_sources") or []
            if not sources:
                issues.append(
                    f"Run '{label}': consolidated entry has empty _consolidate_sources."
                )
        elif not run.get("file"):
            issues.append(f"Run '{label}': missing 'file' key.")
        elif not is_consolidated:
            file_path = config.root_folder / run["file"]
            if not file_path.exists():
                issues.append(f"Run '{label}': file not found -> {file_path}")
        if not run.get("color"):
            pass
        run_type = run.get("type")
        if run_type and run_type not in _VALID_RUN_TYPES:
            issues.append(
                f"Run '{label}': unknown type '{run_type}'. "
                f"Expected one of: {', '.join(sorted(_VALID_RUN_TYPES))}"
            )
    if config.powerpoint_template and not config.powerpoint_template.exists():
        issues.append(f"PowerPoint template not found: {config.powerpoint_template}")
    return issues

def validate_export_map(plot_definitions: tuple, export_map: Optional[dict]) -> list[str]:
    if not export_map or not plot_definitions:
        return []
    from .plot_definitions import PLOT_TYPE_ORDER as type_prefixes
    generated_names: set[str] = set()
    for group_idx, group in enumerate(plot_definitions):
        if not group:
            continue
        prefix = type_prefixes[group_idx] if group_idx < len(type_prefixes) else "plot"
        for plot_def in group:
            plot_name = getattr(plot_def, "name", None)
            if not plot_name:
                continue
            safe = (
                plot_name.lower()
                .replace(" ", "_")
                .replace("(", "").replace(")", "")
                .replace("/", "_").replace("\\", "_")
            )
            generated_names.add(f"{prefix}_{safe}.png")
    mapped_names: set[str] = set()
    for slide_config in export_map.values():
        for img in slide_config.get("images", []):
            mapped_names.add(img)
    orphans = generated_names - mapped_names
    warnings = []
    if orphans:
        warnings.append(
            f"[WARNING] {len(orphans)} plot(s) generated but NOT in POWERPOINT_EXPORT_MAP:"
        )
        for img in sorted(orphans):
            warnings.append(f"  • {img}")
    return warnings

def run_workflow(
    workflow: str,
    *,
    title: str,
    runs: list,
    waveforms=None,
    scatters=None,
    psds=None,
    histograms=None,
    bars=None,
    boxes=None,
    heatmaps=None,
    powerpoint_template=None,
    powerpoint_output=None,
    export_map=None,
    powerpoint_exports=None,
    vibrations_fit=None,
    fig_size=None,
    cli_description: Optional[str] = None,
    **overrides,
):
    plot_definitions = build_plot_groups(
        waveforms=waveforms, scatters=scatters, psds=psds,
        histograms=histograms, bars=bars, boxes=boxes, heatmaps=heatmaps,
    )
    config = workflow_config(
        workflow,
        title=title,
        runs=runs,
        plot_definitions=plot_definitions,
        powerpoint_template=powerpoint_template,
        powerpoint_output=powerpoint_output,
        export_map=export_map,
        powerpoint_exports=powerpoint_exports,
        vibrations_fit=vibrations_fit,
        fig_size=fig_size,
        **overrides,
    )
    return run_from_config(config, parse_plot_cli(cli_description or title))

def run_from_config(config: PlotJobConfig, cli_args=None):
    configure_logging(verbose=config.verbose)
    resolved_export_map = _resolve_export_map(
        config.export_map, config.plot_definitions, start_slide=config.powerpoint_start_slide,
    )
    resolved_exports: list = []
    if config.powerpoint_exports:
        for i, entry in enumerate(config.powerpoint_exports):
            if not isinstance(entry, (list, tuple)) or len(entry) < 3:
                log.warning(
                    "powerpoint_exports[%d] must be (template, output, export_map[, start_slide]); got %r",
                    i, entry,
                )
                continue
            tpl, out, exp_map = entry[0], entry[1], entry[2]
            start = int(entry[3]) if len(entry) >= 4 else 1
            resolved_exports.append(
                (tpl, out, _resolve_export_map(exp_map, config.plot_definitions, start_slide=start))
            )
    if cli_args is not None and getattr(cli_args, "list_plots", False):
        _print_plot_list(config.plot_definitions)
        return
    runs = _expand_folder_runs(config.runs, config.root_folder)
    config.runs = runs
    if cli_args is not None and getattr(cli_args, "runs", None):
        requested = {r.lower() for r in cli_args.runs}
        runs = [r for r in runs if r["name"].lower() in requested]
        if not runs:
            log.error("No runs matched: %s", cli_args.runs)
            print(f"  Available: {[r['name'] for r in config.runs]}")
            raise SystemExit(1)
    issues = validate_config(config)
    if issues:
        print("\n[ERROR] Configuration validation failed:")
        for issue in issues:
            print(f"  X  {issue}")
        raise SystemExit(1)
    if resolved_export_map:
        orphan_warnings = validate_export_map(config.plot_definitions, resolved_export_map)
        if orphan_warnings:
            for line in orphan_warnings:
                print(line)
            print()
    if cli_args is not None and getattr(cli_args, "dry_run", False):
        _print_dry_run(config, runs, resolved_export_map)
        return
    plot_names = None
    open_output = config.open_output
    if cli_args is not None:
        if getattr(cli_args, "only", None):
            plot_names = cli_args.only
        if getattr(cli_args, "no_open", False):
            open_output = False
    fig_size = config.fig_size
    plot_aspect_ratios = {}
    if config.powerpoint_template and resolved_export_map:
        cache_key = (
            str(config.powerpoint_template),
            json.dumps(resolved_export_map, sort_keys=True),
        )
        plot_aspect_ratios = _ASPECT_RATIO_CACHE.get(cache_key)
        if plot_aspect_ratios is None:
            plot_aspect_ratios = get_template_plot_aspect_ratios(config.powerpoint_template, resolved_export_map)
            _ASPECT_RATIO_CACHE[cache_key] = plot_aspect_ratios
    plotter = DataPlotter(
        root_folder=config.root_folder,
        runs=runs,
        plot_definitions=config.plot_definitions,
        channel_mappings=config.channel_mappings,
        channel_transforms=config.channel_transforms,
        calculated_channels=config.calculated_channels,
        filters=config.filters,
        fig_size=fig_size or DEFAULT_FIG_SIZE,
        units_map=config.units_map,
        plot_aspect_ratios=plot_aspect_ratios,
        scatter_max_points=config.scatter_max_points,
        bar_secondary_axis_ratio=config.bar_secondary_axis_ratio,
        box_plot_settings=config.box_plot_settings,
        output_dir=config.output_dir,
        verbose=config.verbose,
        output_dpi=config.output_dpi,
        resample_rate=config.resample_rate,
        vibrations_fit=config.vibrations_fit,
    )
    if cli_args is not None and getattr(cli_args, "list_channels", False):
        _print_run_channels(plotter.run_data)
        return
    _enforce_channel_typo_check(config.plot_definitions, plotter.run_data)
    if cli_args is not None and getattr(cli_args, "check_only", False):
        from .data_quality_report import (
            build_quality_sections, write_data_quality_report, print_quality_summary,
        )
        sections = build_quality_sections(
            runs, plotter.run_data, config.plot_definitions,
            run_sample_rates=plotter.run_sample_rates,
            outlier_log=plotter._outlier_log,
        )
        report_path = write_data_quality_report(plotter.plots_dir, sections)
        print_quality_summary(sections)
        print(f"\nFull report: {report_path}")
        return
    run_plot_job(
        title=config.title,
        plotter=plotter,
        plot_types=None,
        plot_names=plot_names,
        powerpoint_template=config.powerpoint_template,
        powerpoint_output=config.powerpoint_output,
        export_map=resolved_export_map,
        powerpoint_exports=resolved_exports or None,
        open_output=open_output,
    )
    return plotter

def parse_plot_cli(description: str = "Run plotting job"):
    parser = argparse.ArgumentParser(description=description)
    parser.add_argument(
        "--only", nargs="+", metavar="NAME",
        help="Generate only plots whose name matches (case-insensitive).",
    )
    parser.add_argument(
        "--runs", nargs="+", metavar="RUN",
        help="Process only these runs by name (case-insensitive).",
    )
    parser.add_argument(
        "--no-open", action="store_true", default=False,
        help="Do not auto-open the output folder after completion.",
    )
    parser.add_argument(
        "--list-plots", action="store_true", default=False,
        help="Print all configured plot names and exit.",
    )
    parser.add_argument(
        "--list-channels", action="store_true", default=False,
        help="Load each run and print its available channel names, then exit.",
    )
    parser.add_argument(
        "--check-only", action="store_true", default=False,
        help="Load data, run quality checks, and exit without generating plots.",
    )
    parser.add_argument(
        "--dry-run", action="store_true", default=False,
        help="Preview what would be generated without running the pipeline.",
    )
    return parser.parse_args()

def _print_run_channels(run_data):
    if not run_data:
        print("\n(No runs loaded.)\n")
        return
    per_run = {name: set(df.columns) for name, df in run_data.items()}
    common = set.intersection(*per_run.values()) if per_run else set()
    print("\nAvailable channels")
    print("-" * 50)
    print(f"\n  COMMON to all {len(per_run)} run(s) — {len(common)} channel(s):")
    for ch in sorted(common):
        print(f"    {ch}")
    for name, cols in per_run.items():
        extras = sorted(cols - common)
        if extras:
            print(f"\n  ONLY in '{name}' — {len(extras)} channel(s):")
            for ch in extras:
                print(f"    {ch}")
    print()

def _enforce_channel_typo_check(plot_definitions, run_data):
    if not run_data:
        return
    from .dataplotter import collect_referenced_channels
    import difflib
    referenced = set(collect_referenced_channels(plot_definitions))
    union = set()
    for df in run_data.values():
        union.update(df.columns)
    bogus = sorted(ch for ch in referenced if ch not in union)
    if not bogus:
        return
    lower_to_actual = {c.lower(): c for c in union}
    print("\n[ERROR] Plot definitions reference channels that exist in no loaded run:")
    for ch in bogus:
        cands = difflib.get_close_matches(ch, union, n=3, cutoff=0.6)
        if not cands:
            cands_lower = difflib.get_close_matches(ch.lower(), lower_to_actual.keys(), n=3, cutoff=0.55)
            cands = [lower_to_actual[c] for c in cands_lower]
        hint = f"  did you mean: {', '.join(cands)}?" if cands else ""
        print(f"  X  '{ch}'{hint}")
    print("\nRun with --list-channels to see all available channel names.\n")
    raise SystemExit(1)

def _print_plot_list(plot_definitions):
    from .plot_definitions import PLOT_TYPE_ORDER as type_names
    print("\nConfigured Plots:")
    print("-" * 50)
    total = 0
    for i, group in enumerate(plot_definitions or []):
        if not group:
            continue
        label = type_names[i] if i < len(type_names) else f"group_{i}"
        print(f"\n  {label.upper()} ({len(group)}):")
        for plot_def in group:
            name = getattr(plot_def, "name", plot_def)
            try:
                refs = _plot_referenced_channels(plot_def)
            except Exception:
                refs = []
            if refs:
                shown = ", ".join(refs[:8])
                more = f" (+{len(refs) - 8} more)" if len(refs) > 8 else ""
                print(f"    *  {name}  [{shown}{more}]")
            else:
                print(f"    *  {name}")
            total += 1
    print(f"\n  Total: {total} plot(s)")

def _print_dry_run(config, runs, export_map):
    from .plot_definitions import PLOT_TYPE_ORDER as type_names
    print("\n" + "=" * 60)
    print(f"{'DRY RUN':^60}")
    print("=" * 60)
    print(f"\n  Title:   {config.title}")
    print(f"  Output:  {config.output_dir}")
    print(f"\n  Runs ({len(runs)}):")
    for run in runs:
        print(f"    *  {run['name']} ({run.get('type', '?')}) -- {run.get('file', '?')}")
    available_by_run = {}
    for run in runs:
        file_path = config.root_folder / run.get("file", "")
        cols = set()
        try:
            if file_path.suffix.lower() in (".parquet", ".pq"):
                try:
                    import pyarrow.parquet as pq
                    cols = set(pq.read_schema(str(file_path)).names)
                except Exception:
                    try:
                        from fastparquet import ParquetFile  # type: ignore[import-not-found]
                        cols = set(ParquetFile(str(file_path)).columns)
                    except Exception:
                        cols = set()
            elif file_path.suffix.lower() in (".txt", ".csv"):
                with open(file_path, "r", encoding="utf-8", errors="ignore") as fh:
                    fh.readline()
                    header_line = fh.readline()
                    if header_line:
                        cols = {c.strip() for c in header_line.replace("\t", ",").split(",") if c.strip()}
        except Exception:
            cols = set()
        available_by_run[run["name"]] = cols
    total = 0
    print("\n  Plots:")
    est_bytes = 0
    figsize_default = (10, 8)
    dpi = config.output_dpi if hasattr(config, "output_dpi") else 300
    for i, group in enumerate(config.plot_definitions or []):
        if not group:
            continue
        label = type_names[i] if i < len(type_names) else f"group_{i}"
        print(f"    {label}: {len(group)}")
        for plot_def in group:
            name = getattr(plot_def, "name", str(plot_def))
            refs = _plot_referenced_channels(plot_def)
            missing_per_run = []
            for run in runs:
                avail = available_by_run.get(run["name"]) or set()
                if not avail:
                    continue
                missing = [c for c in refs if c not in avail]
                if missing:
                    missing_per_run.append(f"{run['name']}: {', '.join(missing[:5])}{' ...' if len(missing) > 5 else ''}")
            extra = f"  [missing in: {' | '.join(missing_per_run)}]" if missing_per_run else ""
            print(f"      *  {name}{extra}")
            total += 1
            est_bytes += int(figsize_default[0] * figsize_default[1] * dpi * dpi * 3 / 6)
    print(f"    ---------")
    print(f"    total: {total}")
    if total:
        print(f"    estimated on-disk size: ~{est_bytes/1024/1024:.0f} MB")
    if config.powerpoint_template:
        print(f"\n  PowerPoint: {config.powerpoint_output}")
        if export_map:
            print(f"    Slides mapped: {len(export_map)}")
    print("\n" + "=" * 60 + "\n")

def _plot_referenced_channels(plot_def) -> list:
    out: list = []
    kind = getattr(plot_def, "kind", None)
    if kind == "waveform":
        def _walk(item):
            if isinstance(item, str):
                out.append(item)
            elif isinstance(item, (list, tuple)):
                for v in item:
                    _walk(v)
        _walk(plot_def.channels)
        out.append(plot_def.x_channel)
    elif kind == "scatter":
        out.extend([plot_def.x_channel, plot_def.y_channel])
    elif kind in ("psd", "histogram"):
        ch = plot_def.channel
        if isinstance(ch, (list, tuple)):
            out.extend(ch)
        else:
            out.append(ch)
    elif kind == "bar":
        for m in plot_def.metrics or ():
            if isinstance(m, str):
                out.append(m)
            elif isinstance(m, (list, tuple)) and m:
                out.append(m[0])
    elif kind == "box":
        for ch in plot_def.channels:
            out.append(ch)
    elif kind == "heatmap":
        out.extend([plot_def.x_channel, plot_def.y_channel])
        if plot_def.z_channel:
            out.append(plot_def.z_channel)
    return sorted(set(c for c in out if c))

def build_plot_groups(
    *,
    waveforms=None,
    scatters=None,
    psds=None,
    histograms=None,
    bars=None,
    boxes=None,
    heatmaps=None,
):
    boxes = _expand_box_grids(boxes) if boxes else []
    return tuple(
        group or []
        for group in (waveforms, scatters, psds, histograms, bars, boxes, heatmaps)
    )

def _expand_box_grids(boxes):
    from .plot_definitions import BoxPlotGrid
    expanded = []
    for item in boxes:
        if isinstance(item, BoxPlotGrid) and item.render_mode == "expand":
            expanded.extend(item.expand())
        else:
            expanded.append(item)
    return expanded

def workflow_config(
    workflow: str,
    *,
    title: str,
    runs: list,
    plot_definitions: tuple,
    powerpoint_template=None,
    powerpoint_output=None,
    export_map=None,
    powerpoint_exports=None,
    vibrations_fit=None,
    fig_size=None,
    **overrides,
) -> PlotJobConfig:
    from channel_config import (
        CHANNEL_MAPPINGS, UNITS_MAP, CHANNEL_TRANSFORMS,
        SCATTER_MAX_POINTS, BAR_SECONDARY_AXIS_RATIO, BOX_PLOT_SETTINGS,
    )
    _WORKFLOW_MAP = {
        "correlation": ("CORRELATION_INPUT_DIR", "CORRELATION_OUTPUT_DIR",
                        "CORRELATION_CALCULATED", "CORRELATION_FILTERS"),
        "boxplots":    ("BOXPLOT_INPUT_DIR",     "BOXPLOT_OUTPUT_DIR",
                        "BOXPLOT_CALCULATED",    "BOXPLOT_FILTERS"),
        "dampers":     ("DAMPER_INPUT_DIR",      "DAMPER_OUTPUT_DIR",
                        "DAMPER_CALCULATED",     "DAMPER_FILTERS"),
        "ride_dil":    ("RIDE_DIL_INPUT_DIR",    "RIDE_DIL_OUTPUT_DIR",
                        "RIDE_DIL_CALCULATED",   "RIDE_DIL_FILTERS"),
    }
    import channel_config as _cc
    explicit_root = overrides.pop("root_folder", None)
    explicit_out = overrides.pop("output_dir", None)
    resample_rate = overrides.pop("resample_rate", getattr(_cc, "RESAMPLE_RATE", None))
    if workflow in _WORKFLOW_MAP:
        input_dir, output_dir, calc_attr, filt_attr = _WORKFLOW_MAP[workflow]
        root_folder = explicit_root or getattr(_cc, input_dir)
        out_folder = explicit_out or getattr(_cc, output_dir)
        calculated = overrides.pop("calculated_channels", getattr(_cc, calc_attr))
        filters = overrides.pop("filters", getattr(_cc, filt_attr))
    else:
        from channel_config import get_workflow_dirs
        _root, _out = get_workflow_dirs(workflow)
        root_folder = explicit_root or _root
        out_folder = explicit_out or _out
        calculated = overrides.pop("calculated_channels", getattr(_cc, "CALCULATED_CHANNELS", {}))
        filters = overrides.pop("filters", getattr(_cc, "DEFAULT_FILTERS", {}))
    return PlotJobConfig(
        title=title,
        root_folder=root_folder,
        output_dir=out_folder,
        runs=runs,
        plot_definitions=plot_definitions,
        channel_mappings=overrides.pop("channel_mappings", CHANNEL_MAPPINGS),
        channel_transforms=overrides.pop("channel_transforms", CHANNEL_TRANSFORMS),
        calculated_channels=calculated,
        filters=filters,
        resample_rate=resample_rate,
        units_map=overrides.pop("units_map", UNITS_MAP),
        scatter_max_points=overrides.pop("scatter_max_points", SCATTER_MAX_POINTS),
        bar_secondary_axis_ratio=overrides.pop("bar_secondary_axis_ratio", BAR_SECONDARY_AXIS_RATIO),
        box_plot_settings=overrides.pop("box_plot_settings", BOX_PLOT_SETTINGS),
        fig_size=fig_size,
        powerpoint_template=powerpoint_template,
        powerpoint_output=powerpoint_output,
        export_map=export_map,
        powerpoint_exports=powerpoint_exports,
        vibrations_fit=vibrations_fit,
        **overrides,
    )

def run_plot_job(
    *,
    title,
    plotter,
    plot_types=None,
    plot_names=None,
    powerpoint_template=None,
    powerpoint_output=None,
    export_map=None,
    powerpoint_exports=None,
    open_output=True,
):
    import time as _time
    print("\n" + "=" * 80)
    print(f"{title:^80}")
    print("=" * 80 + "\n")
    from .data_quality_report import (
        build_quality_sections, write_data_quality_report, print_quality_summary,
    )
    sections = build_quality_sections(
        plotter.runs, plotter.run_data, plotter.PLOT_DEFINITIONS,
        run_sample_rates=getattr(plotter, "run_sample_rates", None),
        outlier_log=getattr(plotter, "_outlier_log", None),
    )
    report_path = write_data_quality_report(plotter.plots_dir, sections)
    has_issues = any(values for _, values in sections)
    if has_issues:
        print_quality_summary(sections)
        print(f"  Full report: {report_path}\n")
    pre_mtimes = {
        p: p.stat().st_mtime for p in plotter.plots_dir.rglob("*.png")
    }
    t0 = _time.perf_counter()
    print("\nGenerating plots...")
    plotter.plot_data(plot_types=plot_types, plot_names=plot_names)
    elapsed = _time.perf_counter() - t0
    plot_count = sum(
        1 for p in plotter.plots_dir.rglob("*.png")
        if p not in pre_mtimes or p.stat().st_mtime > pre_mtimes[p]
    )
    print(f"\nGenerated {plot_count} plot(s) in {elapsed:.1f}s -> {plotter.plots_dir}")
    exports: list = []
    if powerpoint_exports:
        exports.extend(powerpoint_exports)
    if powerpoint_template and powerpoint_output and export_map:
        exports.append((powerpoint_template, powerpoint_output, export_map))
    for i, exp in enumerate(exports):
        if len(exp) == 3:
            tpl, out, exp_map = exp
        elif len(exp) == 4:
            tpl, out, exp_map, _start = exp
        else:
            log.warning("PowerPoint export entry #%d has invalid shape: %r", i, exp)
            continue
        if not (tpl and out and exp_map):
            continue
        label = Path(out).name
        print(f"\nExporting to PowerPoint [{i+1}/{len(exports)}]: {label}")
        try:
            Path(out).parent.mkdir(parents=True, exist_ok=True)
            export_report_to_powerpoint(
                template_path=tpl,
                output_path=out,
                plots_dir=plotter.plots_dir,
                export_map=exp_map,
                visible=False,
            )
            try:
                os.startfile(out)
            except Exception as open_err:
                log.warning("Could not auto-open PowerPoint file: %s", open_err)
                print(f"File saved to: {out}")
        except Exception as export_err:
            log.error("PowerPoint export failed (%s): %s", label, export_err)
            traceback.print_exc()
    if open_output:
        try:
            os.startfile(plotter.plots_dir)
        except Exception:
            pass
    print("\n" + "=" * 80)
    print(f"{'PROCESSING COMPLETE':^80}")
    print("=" * 80 + "\n")

from .plot_definitions import (  # noqa: E402
    Marker,
    WaveformPlot,
    ScatterPlot,
    PsdPlot,
    HistogramPlot,
    BarPlot,
    BoxPlot,
    BoxPlotGrid,
    HeatmapPlot,
)

MAIN_PLOT_BOX = {
    "left_ratio": 0.079,
    "top_ratio": 0.260,
    "width_ratio": 0.90,
    "height_ratio": 0.65,
}

DOUBLE_PLOT_LAYOUT = {
    "left_ratio": 0.0,
    "top_ratio": 0.245,
    "width_ratio": 1.2,
    "height_ratio": 0.9,
    "gap_ratio": 0.0,
}

MSO_PICTURE_TYPES = {11, 13}

PPTX_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

def _resolve_box(layout_name, slide_width, slide_height, slot_index=0, slot_count=1):
    if layout_name == "main_plot":
        box = MAIN_PLOT_BOX
        return (
            slide_width * box["left_ratio"],
            slide_height * box["top_ratio"],
            slide_width * box["width_ratio"],
            slide_height * box["height_ratio"],
        )
    if layout_name == "double_plot":
        box = DOUBLE_PLOT_LAYOUT
        left = slide_width * box["left_ratio"]
        top = slide_height * box["top_ratio"]
        width = slide_width * box["width_ratio"]
        height = slide_height * box["height_ratio"]
        gap = slide_width * box["gap_ratio"]
        slot_width = (width - gap) / 2
        return (
            left + slot_index * (slot_width + gap),
            top,
            slot_width,
            height,
        )
    raise ValueError(f"Unsupported PowerPoint layout: {layout_name}")

def _replace_slide_pictures(slide):
    for idx in range(slide.Shapes.Count, 0, -1):
        shape = slide.Shapes(idx)
        if shape.Type in MSO_PICTURE_TYPES:
            shape.Delete()

def _get_picture_boxes(slide):
    boxes = []
    for idx in range(1, slide.Shapes.Count + 1):
        sh = slide.Shapes(idx)
        if sh.Type in MSO_PICTURE_TYPES:
            boxes.append((sh.Left, sh.Top, sh.Width, sh.Height))
    return sorted(boxes, key=lambda b: (b[0], b[1]))

def _get_double_plot_boxes(picture_boxes, slide_width, slide_height):
    if len(picture_boxes) >= 2:
        sorted_boxes = sorted(picture_boxes, key=lambda b: b[0])
        return sorted_boxes[:2]
    layout = DOUBLE_PLOT_LAYOUT
    left = slide_width * layout["left_ratio"]
    top = slide_height * layout["top_ratio"]
    total_width = slide_width * layout["width_ratio"]
    total_height = slide_height * layout["height_ratio"]
    gap = slide_width * layout["gap_ratio"]
    slot_width = max((total_width - gap) / 2, 0)
    return [
        (left, top, slot_width, total_height),
        (left + slot_width + gap, top, slot_width, total_height),
    ]

def _get_main_plot_box(picture_boxes, slide_width, slide_height):
    if picture_boxes:
        _, top, _, height = picture_boxes[0]
        return (0, top, slide_width, height)
    box = MAIN_PLOT_BOX
    return (
        slide_width * box["left_ratio"],
        slide_height * box["top_ratio"],
        slide_width * box["width_ratio"],
        slide_height * box["height_ratio"],
    )

def _add_picture_fit(slide, image_path, left, top, width, height, fill_factor=1.0):
    image_path = str(image_path)
    shape = slide.Shapes.AddPicture(image_path, False, True, 0, 0, -1, -1)
    shape.LockAspectRatio = True
    scale = min(width / shape.Width, height / shape.Height)
    scale *= fill_factor
    shape.Width *= scale
    shape.Height *= scale
    shape.Left = left + (width - shape.Width) / 2
    shape.Top = top + (height - shape.Height) / 2
    shape.Line.Visible = True
    shape.Line.ForeColor.RGB = 0
    shape.Line.Weight = 1
    return shape

def get_template_plot_aspect_ratios(template_path, export_map):
    template_path = Path(template_path).resolve()
    if not template_path.exists():
        log.warning("PowerPoint template not found: %s. Using default aspect ratios.", template_path)
        return {}
    aspect_ratios = {}
    try:
        with ZipFile(template_path) as pptx:
            pres_root = ET.fromstring(pptx.read("ppt/presentation.xml"))
            slide_size = pres_root.find("p:sldSz", PPTX_NS)
            slide_width = int(slide_size.attrib.get("cx", 0)) if slide_size is not None else None
            for slide_num, config in export_map.items():
                slide_xml = f"ppt/slides/slide{slide_num}.xml"
                if slide_xml not in pptx.namelist():
                    continue
                root = ET.fromstring(pptx.read(slide_xml))
                picture_boxes = []
                for pic in root.findall(".//p:pic", PPTX_NS):
                    xfrm = pic.find("p:spPr/a:xfrm", PPTX_NS)
                    if xfrm is None:
                        continue
                    ext = xfrm.find("a:ext", PPTX_NS)
                    off = xfrm.find("a:off", PPTX_NS)
                    if ext is None:
                        continue
                    width = int(ext.attrib.get("cx", 0))
                    height = int(ext.attrib.get("cy", 0))
                    left = int(off.attrib.get("x", 0)) if off is not None else 0
                    top = int(off.attrib.get("y", 0)) if off is not None else 0
                    if width > 0 and height > 0:
                        picture_boxes.append((left, top, width, height))
                picture_boxes.sort(key=lambda b: (b[0], b[1]))
                image_files = config.get("images", [])
                slide_aspects = []
                for i, img_file in enumerate(image_files):
                    if i >= len(picture_boxes):
                        break
                    _, _, w, h = picture_boxes[i]
                    if (
                        config.get("layout") == "main_plot"
                        and slide_width is not None
                        and len(image_files) == 1
                    ):
                        w = slide_width
                    slide_aspects.append((img_file, w / h))
                if (
                    config.get("layout") == "double_plot"
                    and len(slide_aspects) == 2
                    and not all(
                        name.startswith(("scatter_", "psd_", "bar_"))
                        for name, _ in slide_aspects
                    )
                ):
                    avg = sum(a for _, a in slide_aspects) / len(slide_aspects)
                    for img, _ in slide_aspects:
                        aspect_ratios[img] = (avg,)
                else:
                    for img, ar in slide_aspects:
                        aspect_ratios[img] = ar
    except Exception as e:
        log.warning("Error reading template aspect ratios: %s. Using default aspect ratios.", e)
        return {}
    return aspect_ratios

def _export_via_pptx(template_path, output_path, plots_dir, export_map):
    from pptx import Presentation
    from pptx.util import Emu
    from PIL import Image as _PILImage
    prs = Presentation(str(template_path))
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    def _iter_pic_elements(slide):
        sp_tree = slide.shapes._spTree
        return list(sp_tree.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}pic")) + \
               list(sp_tree.iter("{http://schemas.openxmlformats.org/presentationml/2006/main}pic"))
    def _picture_boxes(slide):
        boxes = []
        for pic in _iter_pic_elements(slide):
            xfrm = pic.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm")
            if xfrm is None:
                continue
            ext = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ext")
            off = xfrm.find("{http://schemas.openxmlformats.org/drawingml/2006/main}off")
            if ext is None:
                continue
            try:
                w = int(ext.get("cx", 0))
                h = int(ext.get("cy", 0))
                left = int(off.get("x", 0)) if off is not None else 0
                top = int(off.get("y", 0)) if off is not None else 0
            except (TypeError, ValueError):
                continue
            if w > 0 and h > 0:
                boxes.append((left, top, w, h))
        return sorted(boxes, key=lambda b: (b[0], b[1]))
    def _delete_pictures(slide):
        for pic in _iter_pic_elements(slide):
            parent = pic.getparent()
            if parent is not None:
                parent.remove(pic)
    def _add_picture_fit(slide, image_path, left, top, width, height, fill_factor=1.0):
        try:
            with _PILImage.open(image_path) as im:
                img_w, img_h = im.size
        except Exception:
            img_w, img_h = (1, 1)
        if img_w <= 0 or img_h <= 0:
            img_w, img_h = (1, 1)
        scale = min(width / img_w, height / img_h) * fill_factor
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)
        new_l = int(left + (width - new_w) / 2)
        new_t = int(top + (height - new_h) / 2)
        slide.shapes.add_picture(
            str(image_path),
            Emu(new_l), Emu(new_t),
            width=Emu(new_w), height=Emu(new_h),
        )
    def _is_misc_plot(img_file: str) -> bool:
        base = img_file.rsplit("/", 1)[-1].rsplit("\\", 1)[-1]
        return base.startswith(("scatter_", "psd_", "histogram_", "bar_"))
    for slide_num, config in export_map.items():
        idx = int(slide_num) - 1
        if idx < 0 or idx >= len(prs.slides):
            log.warning("Slide %s out of range (template has %d slides). Skipping.",
                        slide_num, len(prs.slides))
            continue
        slide = prs.slides[idx]
        layout = config.get("layout", "main_plot")
        image_files = config.get("images", [])
        if not image_files:
            continue
        existing = _picture_boxes(slide)
        used_template_boxes = False
        if layout == "main_plot" and len(image_files) == 1:
            target_boxes = [_get_main_plot_box(existing, slide_width, slide_height)]
            used_template_boxes = bool(existing)
        elif layout == "double_plot" and len(image_files) == 2:
            target_boxes = _get_double_plot_boxes(existing, slide_width, slide_height)
            used_template_boxes = len(existing) >= 2
        else:
            target_boxes = existing or [
                _resolve_box(layout, slide_width, slide_height,
                             slot_index=i, slot_count=len(image_files))
                for i in range(len(image_files))
            ]
            used_template_boxes = bool(existing)
        _delete_pictures(slide)
        for i, img_file in enumerate(image_files):
            img_path = (Path(plots_dir) / img_file).resolve()
            if not img_path.exists():
                log.warning("Missing plot for slide %s: %s", slide_num, img_file)
                continue
            if i < len(target_boxes):
                left, top, w, h = target_boxes[i]
            else:
                left, top, w, h = _resolve_box(
                    layout, slide_width, slide_height,
                    slot_index=i, slot_count=len(image_files),
                )
            if used_template_boxes:
                fill_factor = 1.0
            elif layout == "double_plot" and _is_misc_plot(img_file):
                fill_factor = 1.2
            else:
                fill_factor = 1.0
            _add_picture_fit(slide, img_path, left, top, w, h, fill_factor)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

def export_report_to_powerpoint(template_path, output_path, plots_dir, export_map, visible=False):
    template_path = Path(template_path).resolve()
    plots_dir = Path(plots_dir).resolve()
    output_path = Path(output_path).resolve()
    if not template_path.exists():
        raise FileNotFoundError(f"PowerPoint template not found: {template_path}")
    try:
        import pptx  # noqa: F401
        import PIL  # noqa: F401
        _export_via_pptx(template_path, output_path, plots_dir, export_map)
        log.info("PowerPoint exported via python-pptx: %s", output_path)
        return
    except ImportError:
        log.debug("python-pptx not available; falling back to win32com.")
    except Exception as exc:
        log.warning(
            "python-pptx export failed (%s); falling back to win32com.", exc,
        )
    try:
        import win32com.client
    except ImportError as exc:
        raise ImportError(
            "Neither python-pptx nor pywin32 is available for PowerPoint export.\n"
            "Install one of:\n"
            "    pip install python-pptx   # cross-platform (preferred)\n"
            "    pip install pywin32       # Windows only"
        ) from exc
    ppt = win32com.client.Dispatch("PowerPoint.Application")
    ppt.Visible = True
    pres = None
    try:
        pres = ppt.Presentations.Open(str(template_path), WithWindow=visible)
        slide_width = pres.PageSetup.SlideWidth
        slide_height = pres.PageSetup.SlideHeight
        for slide_num, cfg in export_map.items():
            slide = pres.Slides(slide_num)
            layout = cfg["layout"]
            image_list = cfg["images"]
            picture_boxes = _get_picture_boxes(slide)
            used_template_boxes = False
            if layout == "main_plot" and len(image_list) == 1:
                target_boxes = [_get_main_plot_box(picture_boxes, slide_width, slide_height)]
                used_template_boxes = bool(picture_boxes)
            elif layout == "double_plot" and len(image_list) == 2:
                target_boxes = _get_double_plot_boxes(picture_boxes, slide_width, slide_height)
                used_template_boxes = len(picture_boxes) >= 2
            else:
                target_boxes = picture_boxes or [
                    _resolve_box(layout, slide_width, slide_height, slot_index=i, slot_count=len(image_list))
                    for i in range(len(image_list))
                ]
                used_template_boxes = bool(picture_boxes)
            _replace_slide_pictures(slide)
            for i, img in enumerate(image_list):
                img_path = plots_dir / img
                if not img_path.exists():
                    log.warning("Missing plot for slide %d: %s", slide_num, img)
                    continue
                if i < len(target_boxes):
                    left, top, width, height = target_boxes[i]
                else:
                    left, top, width, height = _resolve_box(
                        layout,
                        slide_width,
                        slide_height,
                        slot_index=i,
                        slot_count=len(image_list),
                    )
                _basename = img.rsplit("/", 1)[-1].rsplit("\\", 1)[-1]
                if (
                    not used_template_boxes
                    and layout == "double_plot"
                    and _basename.startswith(("scatter_", "psd_", "histogram_", "bar_"))
                ):
                    fill_factor = 1.2
                else:
                    fill_factor = 1.0
                _add_picture_fit(slide, img_path, left, top, width, height, fill_factor)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            pres.SaveAs(str(output_path))
            final = output_path
        except Exception as exc:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            fallback = output_path.with_name(f"{output_path.stem}_{ts}{output_path.suffix}")
            print(
                f"[WARNING][powerpointexporter] Could not save to {output_path} ({exc}). Using fallback: {fallback}"
            )
            pres.SaveAs(str(fallback))
            final = fallback
        print(f"PowerPoint report saved to: {final}")
    except Exception as exc:
        log.error("PowerPoint export failed: %s", exc)
    finally:
        try:
            ppt.Quit()
        except Exception as quit_err:
            log.warning("Error quitting PowerPoint: %s", quit_err)
        pres = None
        ppt = None
